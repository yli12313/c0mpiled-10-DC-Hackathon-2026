from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colors ──────────────────────────────────────────
BG      = RGBColor(0x05, 0x0a, 0x18)
CYAN    = RGBColor(0x00, 0xd4, 0xff)
BLUE    = RGBColor(0x4f, 0x8e, 0xf7)
PURPLE  = RGBColor(0xa8, 0x55, 0xf7)
GOLD    = RGBColor(0xfd, 0xa0, 0x85)
WHITE   = RGBColor(0xe2, 0xe8, 0xf0)
MUTED   = RGBColor(0x64, 0x74, 0x8b)
SURFACE = RGBColor(0x0a, 0x14, 0x32)
RED     = RGBColor(0xff, 0x4f, 0x6a)
GREEN   = RGBColor(0x00, 0xff, 0x88)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank


def bg(slide):
    """Fill slide background dark."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def box(slide, left, top, width, height,
        fill_color=None, border_color=None, border_width=Pt(1), radius=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.width = border_width
    if border_color:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def label(slide, text, left, top, width, height, color=CYAN, size=Pt(10), align=PP_ALIGN.CENTER, bold=False, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = False
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = size
    run.font.color.rgb = color
    run.font.bold   = bold
    run.font.italic = italic
    run.font.name   = "Arial"
    return txb


def heading(slide, text, left, top, width, height, size=Pt(36), color=CYAN, align=PP_ALIGN.CENTER):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.color.rgb = color
    run.font.bold  = True
    run.font.name  = "Arial Black"
    return txb


def body(slide, lines, left, top, width, height, size=Pt(16), color=WHITE, align=PP_ALIGN.LEFT, spacing=Pt(6)):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = spacing
        run = p.add_run()
        run.text = line
        run.font.size  = size
        run.font.color.rgb = color
        run.font.name  = "Arial"
    return txb


def divider(slide, cx, top, width=Inches(0.8), color=CYAN):
    shape = slide.shapes.add_shape(1, cx - width//2, top, width, Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def card(slide, left, top, width, height):
    s = box(slide, left, top, width, height, fill_color=SURFACE, border_color=CYAN, border_width=Pt(0.75))
    return s


def add_nasa_logo(slide, left, top, size=Inches(1.1)):
    try:
        slide.shapes.add_picture(
            "/Users/yli/Desktop/c0mpiled_10_DC_Hackathon/frontend/nasa-logo.png",
            left, top, size, size
        )
    except Exception:
        pass


# ═══════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)

# Subtle top accent line
s = sl.shapes.add_shape(1, 0, 0, W, Pt(3))
s.fill.solid(); s.fill.fore_color.rgb = CYAN; s.line.fill.background()

# NASA logo centered
add_nasa_logo(sl, W//2 - Inches(0.55), Inches(0.55))

# Event label
label(sl, "C0MPILED_10  ·  DC HACKATHON 2026",
      Inches(1), Inches(1.85), Inches(11.333), Inches(0.35),
      color=CYAN, size=Pt(10))

# Title
heading(sl, "NASA EXOPLANET\nQUERY INTERFACE",
        Inches(1), Inches(2.2), Inches(11.333), Inches(2.2),
        size=Pt(48), color=CYAN)

# Tagline
label(sl, "The universe has answers. Ask them.",
      Inches(1), Inches(4.5), Inches(11.333), Inches(0.5),
      color=GOLD, size=Pt(18), italic=True)

# Subtitle
label(sl, "A natural language AI interface over real NASA exoplanet data —\nbuilt on modern PostgreSQL to replace aging legacy systems.",
      Inches(1.5), Inches(5.1), Inches(10.333), Inches(0.9),
      color=MUTED, size=Pt(14))

# Bottom accent
s = sl.shapes.add_shape(1, 0, H - Pt(3), W, Pt(3))
s.fill.solid(); s.fill.fore_color.rgb = PURPLE; s.line.fill.background()


# ═══════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
s = sl.shapes.add_shape(1, 0, 0, W, Pt(3))
s.fill.solid(); s.fill.fore_color.rgb = CYAN; s.line.fill.background()

label(sl, "THE PROBLEM", Inches(1), Inches(0.4), Inches(11), Inches(0.3),
      color=CYAN, size=Pt(10))

heading(sl, "Legacy Systems Are Holding NASA Back",
        Inches(1), Inches(0.75), Inches(11.333), Inches(1.1),
        size=Pt(34), color=CYAN)

divider(sl, W//2, Inches(1.85))

bullets = [
    ("🛢️  Siloed Data",
     "NASA's scientific data lives in aging Microsoft SQL Server systems —\n"
     "inaccessible to modern AI tooling and downstream applications."),
    ("🔒  Expert-Only Access",
     "Researchers and the public can't query the data without deep SQL\n"
     "expertise. Valuable insights go undiscovered every day."),
    ("🚫  Blocked AI Pipeline",
     "Legacy infrastructure prevents machine learning, open data APIs,\n"
     "real-time dashboards, and LLM integrations from being built."),
    ("⏳  Stalled Migration",
     "Modernization efforts have stalled for years — lacking a compelling,\n"
     "working proof of concept to justify the investment."),
]

for i, (title, desc) in enumerate(bullets):
    row = i // 2
    col = i % 2
    cx = Inches(0.5) + col * Inches(6.45)
    cy = Inches(2.1) + row * Inches(2.35)
    card(sl, cx, cy, Inches(6.2), Inches(2.1))
    label(sl, title, cx + Inches(0.2), cy + Inches(0.18),
          Inches(5.8), Inches(0.35), color=CYAN, size=Pt(12), bold=True, align=PP_ALIGN.LEFT)
    body(sl, [desc], cx + Inches(0.2), cy + Inches(0.55),
         Inches(5.8), Inches(1.3), size=Pt(13), color=WHITE, spacing=Pt(0))


# ═══════════════════════════════════════════════════
# SLIDE 3 — The Solution
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
s = sl.shapes.add_shape(1, 0, 0, W, Pt(3))
s.fill.solid(); s.fill.fore_color.rgb = CYAN; s.line.fill.background()

label(sl, "THE SOLUTION", Inches(1), Inches(0.4), Inches(11), Inches(0.3),
      color=CYAN, size=Pt(10))

heading(sl, "Modernize. Query. Discover.",
        Inches(1), Inches(0.75), Inches(11.333), Inches(0.8),
        size=Pt(34), color=CYAN)

divider(sl, W//2, Inches(1.6))

# Flow steps
flow_steps = [
    ("🛰️", "NASA Data"),
    ("→",  ""),
    ("🐘", "PostgreSQL"),
    ("→",  ""),
    ("🤖", "LLM (Groq)"),
    ("→",  ""),
    ("💬", "Plain English"),
]

step_w = Inches(1.55)
arrow_w = Inches(0.45)
total_w = 4 * step_w + 3 * arrow_w
start_x = (W - total_w) // 2
cy = Inches(1.85)
x = start_x

for icon, lbl in flow_steps:
    if icon == "→":
        label(sl, "→", x, cy + Inches(0.35), arrow_w, Inches(0.5),
              color=CYAN, size=Pt(22), align=PP_ALIGN.CENTER)
        x += arrow_w
    else:
        card(sl, x, cy, step_w, Inches(1.2))
        label(sl, icon, x, cy + Inches(0.08), step_w, Inches(0.45),
              color=WHITE, size=Pt(24), align=PP_ALIGN.CENTER)
        label(sl, lbl, x, cy + Inches(0.55), step_w, Inches(0.5),
              color=CYAN, size=Pt(10), align=PP_ALIGN.CENTER)
        x += step_w

# Two cards below
cards = [
    ("🔄  Legacy Migration",
     "Automated pipeline pulls NASA Exoplanet Archive data into PostgreSQL — "
     "structured, indexed, and ready for AI queries from day one."),
    ("✨  Natural Language Query",
     "Ask anything in plain English. The LLM writes SQL, executes it, "
     "and explains results in a vivid 2-3 sentence answer — no expertise needed."),
]

for i, (title, desc) in enumerate(cards):
    cx = Inches(0.5) + i * Inches(6.45)
    cy2 = Inches(3.3)
    card(sl, cx, cy2, Inches(6.2), Inches(3.5))
    label(sl, title, cx + Inches(0.22), cy2 + Inches(0.22),
          Inches(5.8), Inches(0.4), color=CYAN, size=Pt(13), bold=True, align=PP_ALIGN.LEFT)
    body(sl, [desc], cx + Inches(0.22), cy2 + Inches(0.7),
         Inches(5.75), Inches(2.5), size=Pt(14), color=WHITE, spacing=Pt(0))


# ═══════════════════════════════════════════════════
# SLIDE 4 — The Data
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
s = sl.shapes.add_shape(1, 0, 0, W, Pt(3))
s.fill.solid(); s.fill.fore_color.rgb = CYAN; s.line.fill.background()

label(sl, "REAL NASA DATA  ·  LIVE DATABASE", Inches(1), Inches(0.4),
      Inches(11), Inches(0.3), color=CYAN, size=Pt(10))

heading(sl, "20,000+ Records. 5 Tables. One Question.",
        Inches(1), Inches(0.75), Inches(11.333), Inches(0.8),
        size=Pt(32), color=CYAN)

divider(sl, W//2, Inches(1.6))

# Stats row
stats = [("6,273", "Confirmed Exoplanets"), ("7,927", "TESS Candidates"), ("4,701", "Host Stars"), ("1,473", "Asteroid Events")]
sw = Inches(3.0)
sx = (W - 4*sw - 3*Inches(0.13)) // 2
for i, (num, lbl) in enumerate(stats):
    cx = sx + i*(sw + Inches(0.13))
    card(sl, cx, Inches(1.85), sw, Inches(1.3))
    label(sl, num, cx, Inches(1.95), sw, Inches(0.7),
          color=CYAN, size=Pt(34), bold=True, align=PP_ALIGN.CENTER)
    label(sl, lbl, cx, Inches(2.68), sw, Inches(0.35),
          color=MUTED, size=Pt(10), align=PP_ALIGN.CENTER)

# Example queries panel
card(sl, Inches(0.5), Inches(3.35), Inches(12.333), Inches(3.65))

label(sl, "EXAMPLE QUERIES", Inches(0.75), Inches(3.55),
      Inches(11.8), Inches(0.3), color=CYAN, size=Pt(10), align=PP_ALIGN.LEFT)

examples = [
    ('❝  Which mission found the most planets?',     'Kepler, with 2,784 confirmed discoveries.'),
    ('❝  Find Earth-sized planets in the habitable zone.',  'Results with orbital period, temperature & distance.'),
    ('❝  Which asteroid passed closest to Earth in 2024?',  'Name, miss distance, and velocity in seconds.'),
]

for i, (q, a) in enumerate(examples):
    cy = Inches(3.95) + i * Inches(0.95)
    label(sl, q, Inches(0.75), cy, Inches(7.5), Inches(0.4),
          color=WHITE, size=Pt(13), align=PP_ALIGN.LEFT, italic=True)
    label(sl, "→  " + a, Inches(0.75), cy + Inches(0.38), Inches(11.5), Inches(0.38),
          color=CYAN, size=Pt(12), align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════
# SLIDE 5 — Vision / CTA
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
s = sl.shapes.add_shape(1, 0, 0, W, Pt(3))
s.fill.solid(); s.fill.fore_color.rgb = PURPLE; s.line.fill.background()

add_nasa_logo(sl, W//2 - Inches(0.55), Inches(0.3))

label(sl, "THE VISION", Inches(1), Inches(1.55), Inches(11), Inches(0.3),
      color=CYAN, size=Pt(10))

heading(sl, "A Blueprint for Every Agency",
        Inches(1), Inches(1.85), Inches(11.333), Inches(0.9),
        size=Pt(36), color=CYAN)

divider(sl, W//2, Inches(2.75))

label(sl, "This is more than a demo. It's proof that decades of siloed government data\n"
          "can be unlocked for researchers, policymakers, and the public — through AI.",
      Inches(1.5), Inches(2.9), Inches(10.333), Inches(0.8),
      color=WHITE, size=Pt(15), align=PP_ALIGN.CENTER, italic=True)

vision_cards = [
    ("🚀  Immediate Impact",
     "Researchers and citizens can query 6,273 worlds instantly — no SQL, no middleman, no waiting for a data team."),
    ("🏛️  Scales to Any Agency",
     "The same pipeline applies to IRS tax data, EPA environmental records, or any legacy SQL Server system in government."),
]

for i, (title, desc) in enumerate(vision_cards):
    cx = Inches(0.5) + i * Inches(6.45)
    card(sl, cx, Inches(3.85), Inches(6.2), Inches(2.3))
    label(sl, title, cx + Inches(0.22), cx * 0 + Inches(3.85) + Inches(0.2),
          Inches(5.8), Inches(0.4), color=CYAN, size=Pt(13), bold=True, align=PP_ALIGN.LEFT)
    body(sl, [desc], cx + Inches(0.22), Inches(3.85) + Inches(0.65),
         Inches(5.75), Inches(1.5), size=Pt(13), color=WHITE, spacing=Pt(0))

# Tag pills as text row
tags = "PostgreSQL Migration  ·  Text-to-SQL  ·  Open Government Data  ·  AI Infrastructure  ·  NASA · IRS · EPA"
label(sl, tags, Inches(0.5), Inches(6.35), Inches(12.333), Inches(0.45),
      color=MUTED, size=Pt(11), align=PP_ALIGN.CENTER)

s = sl.shapes.add_shape(1, 0, H - Pt(3), W, Pt(3))
s.fill.solid(); s.fill.fore_color.rgb = CYAN; s.line.fill.background()


# ── Save ─────────────────────────────────────────
out = "/Users/yli/Desktop/c0mpiled_10_DC_Hackathon/NASA_Exoplanet_Pitch.pptx"
prs.save(out)
print(f"Saved: {out}")
